from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_premium_presentation():
    prs = Presentation()

    # --- Design Constants ---
    BLACK_BG = RGBColor(18, 18, 18)    # #121212 Soft Black
    GOLD = RGBColor(212, 175, 55)      # #D4AF37 Metallic Gold
    WHITE = RGBColor(250, 250, 250)    # Off-white
    GRAY = RGBColor(160, 160, 160)     # Muted text
    
    FONT_TITLE = 'Arial' # Ideally Playfair, but safe fallback
    FONT_BODY = 'Arial'  # Ideally Montserrat
    
    # helper to clean slide
    def add_blank_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 is usually blank
        # Set Dark Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BLACK_BG
        return slide

    def add_gold_accents(slide):
        # Top Gold Line
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = GOLD
        shape.line.fill.background() # No border

        # Bottom subtle footer line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(7.2), Inches(9), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GRAY
        line.line.fill.background()

    def add_text(slide, text, top, left, width, height, font_size, color, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT_BODY
        p.alignment = align
        return p

    # --- SLIDES ---

    # 1. INTRO SLIDE (Impactful, Minimal)
    slide = add_blank_slide()
    # Big Gold Box in Center
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2.5), Inches(4), Inches(2.5))
    box.fill.background() # Transparent center
    box.line.color.rgb = GOLD
    box.line.width = Pt(3)
    
    add_text(slide, "TLS REAL ESTATE", Inches(3.2), Inches(0), Inches(10), Inches(1), 40, GOLD, True, PP_ALIGN.CENTER)
    add_text(slide, "ELEVATING THE STANDARD", Inches(4), Inches(0), Inches(10), Inches(1), 16, WHITE, False, PP_ALIGN.CENTER)
    add_text(slide, "A Digital Transformation Proposal", Inches(6.5), Inches(0), Inches(10), Inches(0.5), 12, GRAY, False, PP_ALIGN.CENTER)


    # 2. VISION (The "Why")
    slide = add_blank_slide()
    add_gold_accents(slide)
    add_text(slide, "The New Vision", Inches(0.5), Inches(0.5), Inches(5), Inches(1), 32, GOLD, True)
    
    add_text(slide, "From 'Functional' to 'Unforgettable'", Inches(1.5), Inches(0.5), Inches(9), Inches(1), 24, WHITE)
    add_text(slide, 
             "We aren't just updating a website. We are building a digital flagship store. "
             "Your online presence should be as impeccable as the properties you represent.\n\n"
             "•  Authority\n"
             "•  Trust\n"
             "•  Sophistication", 
             Inches(2.5), Inches(0.5), Inches(9), Inches(4), 18, GRAY)


    # 3. VISUAL IDENTITY (Show, Don't Tell - Color Swatches)
    slide = add_blank_slide()
    add_gold_accents(slide)
    add_text(slide, "Modern Opulence", Inches(0.5), Inches(0.5), Inches(5), Inches(1), 32, GOLD, True)
    add_text(slide, "A color palette defined by confidence.", Inches(1.2), Inches(0.5), Inches(9), Inches(1), 14, GRAY)

    # Swatch 1: Black
    s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(2), Inches(2))
    s1.fill.solid(); s1.fill.fore_color.rgb = BLACK_BG
    s1.line.color.rgb = GRAY
    add_text(slide, "Deep Carbon", Inches(4.7), Inches(1), Inches(2), Inches(0.5), 14, WHITE, True, PP_ALIGN.CENTER)
    
    # Swatch 2: Gold
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.5), Inches(2), Inches(2))
    s2.fill.solid(); s2.fill.fore_color.rgb = GOLD
    s2.line.fill.background()
    add_text(slide, "Metallic Gold", Inches(4.7), Inches(4), Inches(2), Inches(0.5), 14, GOLD, True, PP_ALIGN.CENTER)

    add_text(slide, "Signifies depth, power, and prestige.", Inches(5.5), Inches(0.5), Inches(9), Inches(1), 18, WHITE, False, PP_ALIGN.CENTER)


    # 4. EXPERIENCE (Simple bullet points)
    slide = add_blank_slide()
    add_gold_accents(slide)
    add_text(slide, "Client Experience", Inches(0.5), Inches(0.5), Inches(5), Inches(1), 32, GOLD, True)
    
    features = [
        ("Effortless Navigation", "Clients find what they need instantly, whether buying or leasing."),
        ("Visual Listings", "Properties showcased with cinema-quality grids, not clutter."),
        ("Direct Action", "Forms replacing emails. One click to connect.")
    ]
    
    y = 2.0
    for title, desc in features:
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(y+0.1), Inches(0.5), Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = GOLD; dot.line.fill.background()
        
        # Text
        add_text(slide, title, Inches(y), Inches(0.8), Inches(4), Inches(0.5), 20, WHITE, True)
        add_text(slide, desc, Inches(y+0.4), Inches(0.8), Inches(8), Inches(0.5), 16, GRAY)
        y += 1.3


    # 5. ECOSYSTEM (Partners)
    slide = add_blank_slide()
    add_gold_accents(slide)
    add_text(slide, "Your Professional Network", Inches(0.5), Inches(0.5), Inches(6), Inches(1), 32, GOLD, True)
    add_text(slide, 
             "TLS is more than an agency; it's a hub for property solutions.", 
             Inches(1.5), Inches(0.5), Inches(9), Inches(1), 18, WHITE)
    
    # Draw 4 boxes for partners
    labels = ["Mortgage", "Builders", "Staging", "Legal"]
    x = 1.0
    for label in labels:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(x), Inches(1.8), Inches(2))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(30,30,30)
        card.line.color.rgb = GOLD
        add_text(slide, label, Inches(3.8), Inches(x), Inches(1.8), Inches(1), 14, GOLD, True, PP_ALIGN.CENTER)
        x += 2.0


    # 6. PERFORMANCE (Simplified Tech)
    slide = add_blank_slide()
    add_gold_accents(slide)
    add_text(slide, "Built for Speed", Inches(0.5), Inches(0.5), Inches(5), Inches(1), 32, GOLD, True)
    
    add_text(slide, "Your clients don't like to wait. Neither do we.", Inches(1.5), Inches(0.5), Inches(9), Inches(1), 20, WHITE)
    
    add_text(slide, 
             "•  Instant Loading:  Pages load in under 1 second.\n"
             "•  Mobile Perfect:  Looks stunning on iPhones and iPads.\n"
             "•  Safe & Secure:  Bank-grade reliability without the maintenance.",
             Inches(3), Inches(1), Inches(8), Inches(4), 24, GRAY)


    # 7. GROWTH (Simplified SEO)
    slide = add_blank_slide()
    add_gold_accents(slide)
    add_text(slide, "Getting You Found", Inches(0.5), Inches(0.5), Inches(5), Inches(1), 32, GOLD, True)
    
    add_text(slide, 
             "We've optimized every page to help new clients discover TLS.",
             Inches(1.5), Inches(0.5), Inches(9), Inches(1), 20, WHITE)
    
    # Simple arrow flow
    s1 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(3), Inches(1), Inches(2), Inches(1))
    s1.fill.solid(); s1.fill.fore_color.rgb = GOLD
    add_text(slide, "Optimized", Inches(3.2), Inches(1.2), Inches(1.5), Inches(0.5), 14, BLACK_BG, True, PP_ALIGN.CENTER)
    
    s2 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(3), Inches(3.5), Inches(2), Inches(1))
    s2.fill.solid(); s2.fill.fore_color.rgb = GOLD
    add_text(slide, "Discovered", Inches(3.2), Inches(3.7), Inches(1.5), Inches(0.5), 14, BLACK_BG, True, PP_ALIGN.CENTER)
    
    s3 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(3), Inches(6), Inches(2), Inches(1))
    s3.fill.solid(); s3.fill.fore_color.rgb = GOLD
    add_text(slide, "Connected", Inches(3.2), Inches(6.2), Inches(1.5), Inches(0.5), 14, BLACK_BG, True, PP_ALIGN.CENTER)


    # 8. CONCLUSION
    slide = add_blank_slide()
    add_gold_accents(slide)
    add_text(slide, "The Next Step", Inches(2), Inches(0), Inches(10), Inches(1), 40, GOLD, True, PP_ALIGN.CENTER)
    add_text(slide, "Let's bring this vision to life.", Inches(3.5), Inches(0), Inches(10), Inches(1), 24, WHITE, False, PP_ALIGN.CENTER)
    
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(3.5), Inches(3), Inches(1))
    btn.fill.background()
    btn.line.color.rgb = GOLD
    btn.line.width = Pt(2)
    add_text(slide, "VIEW PROTOTYPE", Inches(5.3), Inches(3.5), Inches(3), Inches(1), 16, GOLD, True, PP_ALIGN.CENTER)


    output_path = "TLS_Premium_Proposal.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_premium_presentation()
